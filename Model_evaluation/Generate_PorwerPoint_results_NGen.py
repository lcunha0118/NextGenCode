#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Nov 19 09:20:43 2021

@author: west
"""
import pandas as pd
import geopandas as gp
from datetime import datetime,timedelta
import os
import matplotlib.pyplot as plt
import glob
import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx import Presentation
import df2img
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.enum.section import WD_SECTION


def change_orientation2():
    current_section = document.sections[-1]
    new_width, new_height = current_section.page_height, current_section.page_width
    # new_section = document.add_section(WD_SECTION.NEW_PAGE)
    current_section.orientation = WD_ORIENT.LANDSCAPE
    current_section.page_width = new_width
    current_section.page_height = new_height

    return current_section

# for example

Hyd_folder="/home/west/Projects/CAMELS/CAMELS_Files_Ngen/"
Output_folder="/media/west/Expansion/Projects/CAMELS/CAMELS_Files_Ngen/Results/"
#Hyd_folder="/home/west/Projects/CAMELS/CAMELS_Files_Ngen//"
Output_ngen=["Output_ngen_03292022"]


#CAMELS_list_516="/home/west/Projects/CAMELS/CAMELS_Files_Ngen//CAMELS_v2_list.csv"
# CAMELS_clim_file="/home/west/Projects/CAMELS/Attributes/camels_attributes_v2.0/camels_clim.txt"
# CAMELS_clim=pd.read_csv(CAMELS_clim_file, sep = ';',dtype={'gauge_id': str,'huc_02': str})
# CAMELS_clim=CAMELS_clim.set_index(['gauge_id']) 
# CAMELS_clim.index.names = ['hru_id_CAMELS']

# CAMELS_list_516="/home/west/Projects/CAMELS/CAMELS_Files_Ngen/CAMELS_v2_selected_withNcat.csv"
# CAMELS_516=pd.read_csv(CAMELS_list_516,dtype={'hru_id_CAMELS': str})
# CAMELS_516=CAMELS_516.set_index(['hru_id_CAMELS'])   

# CAMELS_516=pd.concat([CAMELS_516,CAMELS_clim[['aridity','p_seasonality','frac_snow']]],axis=1,join="inner")


CAMELS_list_516="/home/west/Projects/CAMELS/CAMELS_Files_Ngen/CAMELS_v3_complete.csv"
CAMELS_516=pd.read_csv(CAMELS_list_516,dtype={'hru_id_CAMELS': str})
CAMELS_516=CAMELS_516.set_index(['hru_id_CAMELS'])
#CAMELS_516=CAMELS_516[CAMELS_516['SA_analysis']==1]

CAMELS_516=CAMELS_516[(CAMELS_516.frac_snow<0.1) & (CAMELS_516.NCat<=10)]


General_Dir="/data_CAMELS/"
ID_format="id"
Models=["NOAH_CFE","NOAH_CFE_KNash","NOAH_CFE_klNash","NOAH_CFE_GWIC8","PET_1_CFE","PET_2_CFE","PET_3_CFE","PET_4_CFE","PET_5_CFE","NOAH_CFE_X","PET_1_CFE_X","PET_2_CFE_X","PET_3_CFE_X","PET_4_CFE_X","PET_5_CFE_X","NOAH_Topmodel","PET_1_Topmodel","PET_2_Topmodel","PET_3_Topmodel","PET_4_Topmodel","PET_5_Topmodel"]
Models=["NOAH_CFE","NOAH_CFE_X","NOAH_Topmodel"]
Config_file=[Hyd_folder+"CFE_output_var_config.csv",Hyd_folder+"Topmodel_output_var_config.csv"]

# Models=["NOAH_CFE"]
# Config_file=[Hyd_folder+"CFE_output_var_config.csv"]
Spinnup=365 # in Days
min_date_plot=datetime(2007,10,1,0,0)
max_date_plot=datetime(2013,10,1,0,0)
len(CAMELS_516)

prs = Presentation()
title_slide_layout=prs.slide_layouts[6]


slide=prs.slides.add_slide(title_slide_layout)
output_figure=Hyd_folder+"Selected_sites_Map.png"
pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(0.5),height=Inches(7))  
#CAMELS_516=pd.read_csv("/home/west/Projects/CAMELS/HUC01_camels_calib.txt",header=None,dtype=str)
top=Inches(0.1)
left=Inches(0.1)
height=Inches(0.2)
width=Inches(10)

Powerpoint=Output_folder+"Results_Selected_03292022_Pre.pptx"
a_doc_name=Output_folder+"Appendix_A_05202022.docx"
if(os.path.exists(a_doc_name)): os.remove(a_doc_name)
b_doc_name=Output_folder+"Appendix_B_05202022.docx"
if(os.path.exists(b_doc_name)): os.remove(b_doc_name)
c_doc_name=Output_folder+"Appendix_C_05202022.docx"
if(os.path.exists(c_doc_name)): os.remove(c_doc_name)
d_doc_name=Output_folder+"Appendix_D_05202022.docx"
if(os.path.exists(d_doc_name)): os.remove(d_doc_name)
e_doc_name=Output_folder+"Appendix_E_05202022.docx"
if(os.path.exists(e_doc_name)): os.remove(e_doc_name)

if(os.path.exists(Powerpoint)): os.remove(Powerpoint)

document = Document()
change_orientation2()
documentc = Document()

document.add_heading('Appendix A: Preliminary model evaluation results - hydrographs', 0)
documentc.add_heading('Appendix C: Daily runoff distribution model intercomparison', 0)
p = document.add_paragraph()
p = document.add_paragraph("Formulations:")
p = document.add_paragraph("NOAH_CFE: Noah-OWP-Modular with CFE (Schaake)")
p = document.add_paragraph("NOAH_CFE_Knash: Noah-OWP-Modular with CFE (Schaake) – modified k_Nash parameter from 0.03 to 0.00003")
p = document.add_paragraph("NOAH_CFE_klNash: Noah-OWP-Modular with CFE (Schaake) – modified k_lf parameter from 0.01 to 0.0001")
p = document.add_paragraph("NOAH_CFE_GWIC4: Noah-OWP-Modular with CFE (Schaake)  – modified GW storage initial condition from 0.05 to 0.4")
p = document.add_paragraph("NOAH_CFE_GWIC8: Noah-OWP-Modular with CFE (Schaake) – modified GW storage initial condition from 0.05 to 0.8")
p = document.add_paragraph("PET_(N)_CFE: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = document.add_paragraph("NOAH_CFE_X: Noah-OWP-Modular with CFE (Xinanjiang)")
p = document.add_paragraph("PET_(N)_CFE_X: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = document.add_paragraph("NOAH_Topmodel: Noah-OWP-Modular with Topmodel")
p = document.add_paragraph("PET_(N)_Topmodel: PET with Topmodel – N varies from 1 to 5 and indicates the PET model used")
p = document.add_paragraph("N: (1) energy balance, (2) aerodynamic, (3) combined, (4) Priestley-Taylor, and (5) Penman-Monteith methods")

count=0
count_C=0
for i in range (0,len(CAMELS_516)):         
    hru_id=CAMELS_516.index[i]
    Folder=CAMELS_516.iloc[i]['Folder_CAMELS']
    outfolder=Hyd_folder+"/"+Folder+"/"
    catchment_file=Hyd_folder+"/"+Folder+'/spatial/catchment_data.geojson'                   
    zones = gp.GeoDataFrame.from_file(catchment_file)
    total_area=zones['area_sqkm'].sum()

    slide=prs.slides.add_slide(title_slide_layout)
    # text = "hru_id - " + hru_id+" ("+str(round(total_area,2))+" km2) - All Models: \n"
    # text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
    # text = text+"frac_snow: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
    # text = text+"p_seas: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))  
    text = "Basin: " + hru_id+" , Area: "+str(round(total_area,2))+" km2 \n"
    text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
    text = text+"snow fraction: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
    text = text+"seasonality: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))    
    # if(CAMELS_516.iloc[i]['RunPET']==0):
    #     text = text+"\n - Snow dominated - only run NOAH-OWP"
    txBox=slide.shapes.add_textbox(left,top,width,height)
    tf = txBox.text_frame
    tf.text = text
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    for z in range(0,len(Output_ngen)):
        output_figure=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+"Multiple_Models"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"                         
        if(os.path.exists(output_figure)):         
            print( "found "+output_figure)
            pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(1.5),width=Inches(10)) 
            
            p = documentc.add_paragraph()
            p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
            r = p.add_run()
            r.add_picture(output_figure, width=Inches(6))
            count_C=count_C+1
            p = documentc.add_paragraph("Figure C."+str(count_C)+": "+text.replace("\n"," (").replace(" - ",", ")+")")
            
        for j in range (0,len(Models)): 
            
            if ("PET" in Models[j]) and (CAMELS_516.iloc[i]['RunPET']==0):
                print("Snow dominated areas, do not plot PET")
            else:
                # Read output config file
                

                
                #Obs_Q_file=Hyd_folder+"/"+Folder+"/Validation/usgs_hourly_flow_2007-2019_"+hru_id+".csv"
                #Obs_Q_cms=pd.read_csv(Obs_Q_file,parse_dates=True,index_col=1)
                #Obs_Q_cms=Obs_Q_cms['q_cms']
                #Obs_q_f3s_max_year=(Obs_Q_cms.resample('Y').max()*35.31).to_frame().round(1)
                #Obs_q_f3s_max_year=Obs_q_f3s_max_year.rename(columns={"q_cms":"q_cfs"})
                Results=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+Models[j]+"/"
                title_slide_layout=prs.slide_layouts[6]
                slide=prs.slides.add_slide(title_slide_layout)
                text = "Formulation " + str(Models[j]) + ", Basin: " + hru_id+" , Area: "+str(round(total_area,2))+" km2 \n"
                text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
                text = text+"snow fraction: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
                text = text+"seasonality: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))
                txBox=slide.shapes.add_textbox(left,top,width,height)
                tf = txBox.text_frame
                tf.text = text
                txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                output_figure1= Results+"Runoff_figure"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"            
                if(os.path.exists(output_figure1)):
                    
                    print (hru_id + " " + Models[j])
                    pic=slide.shapes.add_picture(output_figure1,left=Inches(0.1),top=Inches(1.0),height=Inches(5.8))  
                    output_figure2=Results+"Scatter_plot"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"
                    if(os.path.exists(output_figure2)):                                    
                        pic=slide.shapes.add_picture(output_figure2,left=Inches(6.5),top=Inches(1.0),height=Inches(5.8)) 
                        count=count+1
                        p = document.add_paragraph()
                        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
                        r = p.add_run()
                        # r.add_text('')

                        r.add_picture(output_figure1, width=Inches(4.9))
                        r2 = p.add_run()
                        r2.add_picture(output_figure2, width=Inches(2.7))
                        caption=": (a) Cumulative flux (m), (b) Simulated and observed runoff (m/h), (c) All components (m), (d) Storage and deficit (m), (e) Scatter plot of observed and simulated daily runoff versus simulated, and (e) seasonality of observed and simulated daily runoff"
                        p = document.add_paragraph("Figure A."+str(count)+": "+text.replace("\n"," (").replace(" - ",", ")+")"+caption)
                        
                        # fig = df2img.plot_dataframe(
                        #     Obs_q_f3s_max_year,
                        #     title=dict(
                        #         font_color="darkred",
                        #         font_family="Times New Roman",
                        #         font_size=16,
                        #         text="Max annual flow (ft3/s) - " + hru_id,
                        #     ),
                        #     tbl_header=dict(
                        #         align="right",
                        #         fill_color="blue",
                        #         font_color="white",
                        #         font_size=12,
                        #         line_color="darkslategray",
                        #     ),
                        #     tbl_cells=dict(
                        #         align="right",
                        #         line_color="darkslategray",
                        #     ),
                        #     row_fill_color=("#ffffff", "#d7d8d6"),
                        #     fig_size=(300,200),
                        # )
                        # output_figure3= Results+"Table_Max_obs.png"            
                        
                        # df2img.save_dataframe(fig=fig, filename=output_figure3)
                        # pic=slide.shapes.add_picture(output_figure3,left=Inches(0.1),top=Inches(6),height=Inches(1.6)) 
               
prs.save(Powerpoint)     
document.save(a_doc_name)
documentc.save(c_doc_name)


prs = Presentation()
title_slide_layout=prs.slide_layouts[6]
slide=prs.slides.add_slide(title_slide_layout)
output_figure=Output_folder+"Selected_sites_Map.png"
pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(0.5),height=Inches(7))  


documentD = Document()
documentD.add_heading('Appendix D: Hydrographs for WY 2012 with runoff components', 0)

p = documentD.add_paragraph("Formulations:")
p = documentD.add_paragraph("NOAH_CFE: Noah-OWP-Modular with CFE (Schaake)")
p = documentD.add_paragraph("NOAH_CFE_Knash: Noah-OWP-Modular with CFE (Schaake) – modified k_Nash parameter from 0.03 to 0.00003")
p = documentD.add_paragraph("NOAH_CFE_klNash: Noah-OWP-Modular with CFE (Schaake) – modified k_lf parameter from 0.01 to 0.0001")
p = documentD.add_paragraph("NOAH_CFE_GWIC4: Noah-OWP-Modular with CFE (Schaake)  – modified GW storage initial condition from 0.05 to 0.4")
p = documentD.add_paragraph("NOAH_CFE_GWIC8: Noah-OWP-Modular with CFE (Schaake) – modified GW storage initial condition from 0.05 to 0.8")
p = documentD.add_paragraph("PET_(N)_CFE: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = documentD.add_paragraph("NOAH_CFE_X: Noah-OWP-Modular with CFE (Xinanjiang)")
p = documentD.add_paragraph("PET_(N)_CFE_X: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = documentD.add_paragraph("NOAH_Topmodel: Noah-OWP-Modular with Topmodel")
p = documentD.add_paragraph("PET_(N)_Topmodel: PET with Topmodel – N varies from 1 to 5 and indicates the PET model used")
p = documentD.add_paragraph("N: (1) energy balance, (2) aerodynamic, (3) combined, (4) Priestley-Taylor, and (5) Penman-Monteith methods")

count_D=0
Powerpoint=Powerpoint.replace(".pptx","by_year.pptx")
if(os.path.exists(Powerpoint)): os.remove(Powerpoint)
for i in range (0,len(CAMELS_516)):         
    hru_id=CAMELS_516.index[i]
    Folder=CAMELS_516.iloc[i]['Folder_CAMELS']
    outfolder=Output_folder+"/"+Folder+"/"
    catchment_file=Hyd_folder+"/"+Folder+'/spatial/catchment_data.geojson'                   
    zones = gp.GeoDataFrame.from_file(catchment_file)
    total_area=zones['area_sqkm'].sum()

    slide=prs.slides.add_slide(title_slide_layout)
    text = "hru_id - " + hru_id+" ("+str(round(total_area,2))+" km^2) - All Models: \n"
    text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
    text = text+"frac_snow: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
    text = text+"p_seas: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))     
    if(CAMELS_516.iloc[i]['RunPET']==0):
        text = text+"\n - Snow dominated - only run NOAH-OWP"
    txBox=slide.shapes.add_textbox(left,top,width,height)
    tf = txBox.text_frame
    tf.text = text
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    for z in range(0,len(Output_ngen)):
        output_figure=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+"Multiple_Models"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"                         
        if(os.path.exists(output_figure)):         
            print( "found "+output_figure)
            pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(1.5),width=Inches(10)) 
            
        for j in range (0,len(Models)): 
            if ("PET" in Models[j]) and (CAMELS_516.iloc[i]['RunPET']==0):
                print("Snow dominated areas, do not plot PET")
            else:
                Results=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+Models[j]+"/"
                output_figure1= Results+"Runoff_figure"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"            
                if(os.path.exists(output_figure1)):
                    
                    print (hru_id + " " + Models[j])
                    pic=slide.shapes.add_picture(output_figure1,left=Inches(0.1),top=Inches(1.0),height=Inches(5.8))  
                    output_figure2=Results+"Scatter_plot"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"
                    if(os.path.exists(output_figure2)):                                    
                        pic=slide.shapes.add_picture(output_figure2,left=Inches(6.5),top=Inches(1.0),height=Inches(5.8)) 
                                
                        # Read output config file
                        
                        for iyear in range(2007,2013):
                            min_date_plot1=datetime(iyear,10,1,0,0)
                            max_date_plot1=datetime(iyear+1,10,1,0,0)
                            Title_str=hru_id+ " - " +Models[j]+' - ' + min_date_plot1.strftime ('%Y-%m-%d')+' through '+ max_date_plot1.strftime ('%Y-%m-%d')
                            output_figure1= Results+"Runoff_figure"+min_date_plot1.strftime ('%Y-%m')+"_"+max_date_plot1.strftime ('%Y-%m')+".png"                
                            if(os.path.exists(output_figure1)):  
                                print("Add year plot " + str(iyear))
                                title_slide_layout=prs.slide_layouts[6]
                                slide=prs.slides.add_slide(title_slide_layout)
                                text = "Formulation: " + str(Models[j]) + ", Basin: " + hru_id+" , Area: "+str(round(total_area,2))+" km2 \n"
                                text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
                                text = text+"snow fraction: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
                                text = text+"seasonality: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))
                                txBox=slide.shapes.add_textbox(left,top,width,height)
                                tf = txBox.text_frame
                                tf.text = text
                                txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                pic=slide.shapes.add_picture(output_figure1,left=Inches(0.1),top=Inches(1.0),height=Inches(5.8)) 
                        for iyear in range(2011,2012):
                            min_date_plot1=datetime(iyear,10,1,0,0)
                            max_date_plot1=datetime(iyear+1,10,1,0,0)
                            Title_str=hru_id+ " - " +Models[j]+' - ' + min_date_plot1.strftime ('%Y-%m-%d')+' through '+ max_date_plot1.strftime ('%Y-%m-%d')
                            output_figure1= Results+"Runoff_figure"+min_date_plot1.strftime ('%Y-%m')+"_"+max_date_plot1.strftime ('%Y-%m')+".png"                
                            if(os.path.exists(output_figure1)):                      
                        
                                p = documentD.add_paragraph()
                                p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                                r = p.add_run()
                                r.add_picture(output_figure1, width=Inches(5.5))
                                count_D=count_D+1
                                caption=": (a) Cumulative flux (m), (b) Simulated and observed runoff (m/h), (c) Runoff components (m), (d) All components (m), (e) Storage and deficit (m), (e) Scatter plot of observed and simulated daily runoff versus simulated, and (e) seasonality of observed and simulated daily runoff"
                                p = documentD.add_paragraph("Figure D."+str(count_D)+": "+text.replace("\n"," (").replace(" - ",", ")+")"+caption)
                               
prs.save(Powerpoint)   
documentD.save(d_doc_name)      

prs = Presentation()
title_slide_layout=prs.slide_layouts[6]
slide=prs.slides.add_slide(title_slide_layout)
output_figure=Output_folder+"Selected_sites_Map.png"
pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(0.5),height=Inches(7))  


Powerpoint=Powerpoint.replace("by_year.pptx","peak.pptx")
if(os.path.exists(Powerpoint)): os.remove(Powerpoint)
for i in range (0,len(CAMELS_516)):         
    hru_id=CAMELS_516.index[i]
    Folder=CAMELS_516.iloc[i]['Folder_CAMELS']
    outfolder=Output_folder+"/"+Folder+"/"
    catchment_file=Hyd_folder+"/"+Folder+'/spatial/catchment_data.geojson'                   
    zones = gp.GeoDataFrame.from_file(catchment_file)
    total_area=zones['area_sqkm'].sum()

    # slide=prs.slides.add_slide(title_slide_layout)
    # text = "hru_id - " + hru_id+" ("+str(round(total_area,2))+" km^2) - All Models: \n"
    # text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
    # text = text+"frac_snow: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
    # text = text+"p_seas: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))     
    # if(CAMELS_516.iloc[i]['RunPET']==0):
    #     text = text+"\n - Snow dominated - only run NOAH-OWP"
    # txBox=slide.shapes.add_textbox(left,top,width,height)
    # tf = txBox.text_frame
    # tf.text = text
    # txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    for z in range(0,len(Output_ngen)):
        output_figure=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+"Multiple_Models"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"                               
        for j in range (0,len(Models)): 
            if ("PET" in Models[j]) and (CAMELS_516.iloc[i]['RunPET']==0):
                print("Snow dominated areas, do not plot PET")
            else:
                Results=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+Models[j]+"/"
                        
                
                output_figure= Results+"Runoff_figure_Peak.png"
                if(os.path.exists(output_figure)) & (total_area<30):  
                    print("Add year plot " + str(iyear))
                    title_slide_layout=prs.slide_layouts[6]
                    slide=prs.slides.add_slide(title_slide_layout)
                    text = "Formulation: " + str(Models[j]) + ", Basin: " + hru_id+" , Area: "+str(round(total_area,2))+" km2 \n"
                    text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
                    text = text+"snow fraction: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
                    text = text+"seasonality: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))
                    txBox=slide.shapes.add_textbox(left,top,width,height)
                    tf = txBox.text_frame
                    tf.text = text
                    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    output_figure1= Results+"Runoff_figure"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"    
                    pic=slide.shapes.add_picture(output_figure1,left=Inches(0.1),top=Inches(1.0),height=Inches(4.8))
                    pic=slide.shapes.add_picture(output_figure,left=Inches(5.2),top=Inches(1.0),height=Inches(5.5)) 
                else:
                    print ("Did not find file "+ output_figure)
   
prs.save(Powerpoint) 

title_slide_layout=prs.slide_layouts[6]
slide=prs.slides.add_slide(title_slide_layout)
output_figure=Output_folder+"Selected_sites_Map.png"
pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(0.5),height=Inches(7))  

Powerpoint=Powerpoint.replace(".pptx","by_year.pptx")

if(os.path.exists(Powerpoint)): os.remove(Powerpoint)

Model_PET=["NOAH-OWP","NOAH-OWP","NOAH-OWP"]
Model_Runoff=["CFE","CFE_X","Topmodel"]
CriteriaAr=["Selected"]
start_time="2007-10-01 00:00:00"
end_time="2013-10-01 00:00:00"

Var_dictionary={'Topmodel':['Qout','land_surface_water__runoff_mass_flux,land_surface_water__baseflow_volume_flux','land_surface_water__baseflow_volume_flux'],
                 'CFE':['Q_OUT','GIUH_RUNOFF','NASH_LATERAL_RUNOFF','DEEP_GW_TO_CHANNEL_FLUX','ACTUAL_ET','POTENTIAL_ET'],
                 'CFE_X':['Q_OUT','GIUH_RUNOFF','NASH_LATERAL_RUNOFF','DEEP_GW_TO_CHANNEL_FLUX','ACTUAL_ET','POTENTIAL_ET']}
Var_name={'Topmodel':['Qout','DIRECT_RUNOFF','BASEFLOW'],
          'CFE':['Q_OUT','DIRECT_RUNOFF','NASH_LATERAL_RUNOFF','DEEP_GW_TO_CHANNEL_FLUX','ACTUAL_ET','POTENTIAL_ET'],
          'CFE_X':['Q_OUT','DIRECT_RUNOFF','NASH_LATERAL_RUNOFF','DEEP_GW_TO_CHANNEL_FLUX','ACTUAL_ET','POTENTIAL_ET']}
ArrayOfStats=['normalized_nash_sutcliffe','volume_error','peak_error_single']


documentB = Document()
documentB.add_heading('Appendix B: Sensitivity analysis results', 0)
p = documentB.add_paragraph("Local sensitivity analysis was performed for a selection of the parameters in Noah-OWP-Modular, CFE and TOPMODEL.We selected the smallest subbasin for each of the 234 CAMELS basins. We then varied each parameter independently within its valid range as shown in the tables in the parameter data section. Similar to the model evaluation runs described above, we ran the sensitivity analysis between 2007-10-01 and 2013-09-30, discarding the first year as a model spin-up period.")
p = documentB.add_paragraph("To evaluate parameter sensitivity we assessed multiple water balance components: total runoff, direct runoff, lateral runoff, groundwater flow (only for CFE), and potential and actual evapotranspiration. For each parameter (e.g. K_Nash for CFE) and variable (e.g total runoff) variation, the relative difference in normalized Nash Sutcliffe (NNSE), total volume and peak were calculated using a baseline run as reference. ")
CAMELS_Folder="/home/west/Projects/CAMELS/CAMELS_Files_Ngen/"
Spinup= 360 #in days
count_B=0
for crit_id in range(0,len(CriteriaAr)):
    Criteria=CriteriaAr[crit_id]
    for index_im in range(0,len(Model_Runoff)):
        m_runoff=  Model_Runoff[index_im]   
        m_PET=  Model_PET[index_im]           
        Gen_description=m_PET+"_"+m_runoff+ "_"+start_time.replace(' 00:00:00',"")+ "_"+end_time.replace(' 00:00:00',"")+"_Spinup"+str(Spinup)                                
        Btw_site_comparison=CAMELS_Folder+"/BTW_site_comparison/"
        Variables=Var_dictionary[m_runoff]
        Variables_names=Var_name[m_runoff] 
        title=Criteria+"_" + m_runoff + "_"+start_time.replace(' 00:00:00',"")+ "_"+end_time.replace(' 00:00:00',"")+" Spinup"+str(Spinup)+"\n"    
        
        all_results=glob.glob(Btw_site_comparison+"*" +m_runoff+"_*"+"Boxplot_Q.png")
        p = documentB.add_paragraph()
        p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        r = p.add_run()
        r.add_picture(all_results[0], width=Inches(6))
        count_B=count_B+1
        p = documentB.add_paragraph("Figure B."+str(count_B)+": Formulation: "+m_runoff+", Variable: Total runoff, Criteria: All")
        
        for ii in range(0,len(Variables)):       
           print(Variables[ii]) 
           for ij in range(0,len(ArrayOfStats)): 
               print(ArrayOfStats[ij]) 
               all_results=glob.glob(Btw_site_comparison+"*" +m_runoff+"_2007*"+Variables[ii]+"*" + ArrayOfStats[ij]+"*.png")
               if(len(all_results)>0):
                   p = documentB.add_paragraph()
                   p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                   r = p.add_run()
                   r.add_picture(all_results[0], width=Inches(4.2))
                   count_B=count_B+1
                   p = documentB.add_paragraph("Figure B."+str(count_B)+": Formulation: "+m_runoff+", Variable: " + Variables[ii] + ", Criteria: " +ArrayOfStats[ij])
     
documentB.save(b_doc_name)  

Models=["NOAH_CFE","PET_1_CFE","PET_2_CFE","PET_3_CFE","PET_4_CFE","PET_5_CFE","NOAH_CFE_X","PET_1_CFE_X","PET_2_CFE_X","PET_3_CFE_X","PET_4_CFE_X","PET_5_CFE_X","NOAH_Topmodel","PET_1_Topmodel","PET_2_Topmodel","PET_3_Topmodel","PET_4_Topmodel","PET_5_Topmodel"]
documentE = Document()
documentE.add_heading('Appendix E: PET and ET', 0)
p = documentE.add_paragraph()
p = documentE.add_paragraph("Formulations:")
p = documentE.add_paragraph("NOAH_CFE: Noah-OWP-Modular with CFE (Schaake)")
p = documentE.add_paragraph("PET_(N)_CFE: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = documentE.add_paragraph("NOAH_CFE_X: Noah-OWP-Modular with CFE (Xinanjiang)")
p = documentE.add_paragraph("PET_(N)_CFE_X: PET with CFE – N varies from 1 to 5 and indicates the PET model used")
p = documentE.add_paragraph("NOAH_Topmodel: Noah-OWP-Modular with Topmodel")
p = documentE.add_paragraph("PET_(N)_Topmodel: PET with Topmodel – N varies from 1 to 5 and indicates the PET model used")
p = documentE.add_paragraph("N: (1) energy balance, (2) aerodynamic, (3) combined, (4) Priestley-Taylor, and (5) Penman-Monteith methods")

count_E=0

for i in range (0,len(CAMELS_516)):         
    hru_id=CAMELS_516.index[i]
    Folder=CAMELS_516.iloc[i]['Folder_CAMELS']
    outfolder=Output_folder+"/"+Folder+"/"
    catchment_file=Hyd_folder+"/"+Folder+'/spatial/catchment_data.geojson'                   
    zones = gp.GeoDataFrame.from_file(catchment_file)
    total_area=zones['area_sqkm'].sum()

    slide=prs.slides.add_slide(title_slide_layout)
    text = "hru_id - " + hru_id+" ("+str(round(total_area,2))+" km^2) - All Models: \n"
    text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
    text = text+"frac_snow: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
    text = text+"p_seas: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))     
    if(CAMELS_516.iloc[i]['RunPET']==0):
        text = text+"\n - Snow dominated - only run NOAH-OWP"
    txBox=slide.shapes.add_textbox(left,top,width,height)
    tf = txBox.text_frame
    tf.text = text
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
    for z in range(0,len(Output_ngen)):

        for j in range (0,len(Models)): 
            if ("PET" in Models[j]) and (CAMELS_516.iloc[i]['RunPET']==0):
                print("Snow dominated areas, do not plot PET")
            else:
                Results=Output_folder+"/"+Folder+"/"+Output_ngen[z]+"/"+Models[j]+"/"
                output_figure= Results+"PET"+min_date_plot.strftime ('%Y-%m')+"_"+max_date_plot.strftime ('%Y-%m')+".png"
                if(os.path.exists(output_figure)):
 
                    text = "PET and ET for Formulation: " + str(Models[j]) + ", Basin: " + hru_id+" , Area: "+str(round(total_area,2))+" km2 \n"
                    text = text+"aridity: "+str(round(CAMELS_516.iloc[i]['aridity'],2))+" - "
                    text = text+"snow fraction: "+str(round(CAMELS_516.iloc[i]['frac_snow'],2))+" - "
                    text = text+"seasonality: "+str(round(CAMELS_516.iloc[i]['p_seasonality'],2))
                    txBox=slide.shapes.add_textbox(left,top,width,height)
                    tf = txBox.text_frame
                    tf.text = text
                    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    pic=slide.shapes.add_picture(output_figure,left=Inches(0.1),top=Inches(1.0),height=Inches(5.8)) 
        
            
                    p = documentE.add_paragraph()
                    p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                    r = p.add_run()
                    r.add_picture(output_figure, width=Inches(6))
                    count_E=count_E+1
                    caption=": (a) Diurnal cycle of PET and ET (m), (b) Annual cycle of PET and ET"
                    p = documentE.add_paragraph("Figure E."+str(count_E)+": "+text.replace("\n"," (").replace(" - ",", ")+")"+caption)
                   
  
documentE.save(e_doc_name)                      